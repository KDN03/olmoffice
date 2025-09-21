import os
import uuid
import subprocess
import logging
import time
import requests
import shutil
from threading import Thread
from flask import Flask, request, send_file, render_template, jsonify
from werkzeug.utils import secure_filename
import tempfile
import uuid
from pathlib import Path

# Import conversion libraries
from PIL import Image
import pdfkit
import pdfplumber
import pandas as pd
from openpyxl import Workbook
import base64
import io

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s %(levelname)s %(name)s %(message)s'
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'converted'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# App configuration
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-key-change-in-production')

# CloudConvert API configuration (optional)
CLOUDCONVERT_API_KEY = os.environ.get('CLOUDCONVERT_API_KEY', '')

# Expanded file support
ALLOWED_EXTENSIONS = {
    'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx',
    'pdf', 'html', 'htm', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff',
    'txt', 'csv', 'rtf', 'odt', 'ods', 'odp'
}

# All conversion formats from the image
CONVERSION_TYPES = {
    'to_pdf': ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'html', 'htm', 'txt', 'csv'],
    'from_pdf': ['jpg', 'jpeg', 'png', 'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'html', 'txt'],
    'pdf_a': ['pdf']
}

def allowed_file(filename):
    """Check if the file has an allowed extension and is safe."""
    if not filename or '.' not in filename:
        return False
    
    ext = filename.rsplit('.', 1)[-1].lower()
    return ext in ALLOWED_EXTENSIONS and len(filename) <= 255

def cleanup_old_files(folder_path, max_age_hours=24):
    """Remove files older than max_age_hours from the specified folder."""
    try:
        current_time = time.time()
        max_age_seconds = max_age_hours * 3600
        
        if not os.path.exists(folder_path):
            return
            
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            if os.path.isfile(file_path):
                file_age = current_time - os.path.getmtime(file_path)
                if file_age > max_age_seconds:
                    try:
                        os.remove(file_path)
                        logger.info(f"Cleaned up old file: {filename}")
                    except OSError as e:
                        logger.error(f"Failed to remove {file_path}: {e}")
    except Exception as e:
        logger.error(f"Cleanup failed for {folder_path}: {e}")

def start_cleanup_thread():
    """Start background thread for periodic cleanup."""
    def cleanup_worker():
        while True:
            time.sleep(3600)  # Run cleanup every hour
            cleanup_old_files(UPLOAD_FOLDER)
            cleanup_old_files(OUTPUT_FOLDER)
    
    cleanup_thread = Thread(target=cleanup_worker, daemon=True)
    cleanup_thread.start()
    logger.info("Cleanup thread started")

def check_libreoffice_installation():
    """Check if LibreOffice is properly installed and accessible."""
    if os.name == 'nt':  # Windows
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        
        for path in libreoffice_paths:
            if os.path.exists(path):
                logger.info(f"Found LibreOffice at: {path}")
                return path
                
        logger.warning("LibreOffice not found in standard installation paths")
        return None
    else:  # Linux/Mac
        try:
            result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                path = result.stdout.strip()
                logger.info(f"Found LibreOffice at: {path}")
                return 'libreoffice'
        except Exception:
            pass
        
        logger.warning("LibreOffice not found")
        return None

def check_wkhtmltopdf_installation():
    """Check if wkhtmltopdf is installed."""
    if os.name == 'nt':  # Windows
        wkhtmltopdf_paths = [
            r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
            r"C:\Program Files (x86)\wkhtmltopdf\bin\wkhtmltopdf.exe",
        ]
        for path in wkhtmltopdf_paths:
            if os.path.exists(path):
                return pdfkit.configuration(wkhtmltopdf=path)
    else:
        # Check common Linux/Unix paths including cloud environments
        common_paths = [
            '/usr/bin/wkhtmltopdf',
            '/usr/local/bin/wkhtmltopdf',
            '/app/.apt/usr/bin/wkhtmltopdf',  # Render.com with Aptfile
            shutil.which('wkhtmltopdf')  # System PATH
        ]
        
        for path in common_paths:
            if path and os.path.exists(path):
                logger.info(f"Found wkhtmltopdf at: {path}")
                if path != shutil.which('wkhtmltopdf'):
                    return pdfkit.configuration(wkhtmltopdf=path)
                else:
                    return None  # Use default configuration
        
        # Try subprocess as fallback
        try:
            result = subprocess.run(['which', 'wkhtmltopdf'], check=True, capture_output=True, timeout=5)
            if result.returncode == 0:
                path = result.stdout.decode().strip()
                logger.info(f"Found wkhtmltopdf via which: {path}")
                return None  # Use default configuration
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            pass
    
    logger.warning("wkhtmltopdf not found")
    return False

class ConversionEngine:
    """Comprehensive conversion engine with multiple methods"""
    
    def __init__(self):
        self.libreoffice_path = check_libreoffice_installation()
        self.wkhtmltopdf_config = check_wkhtmltopdf_installation()
        self.has_libreoffice = self.libreoffice_path is not None
        self.has_wkhtmltopdf = self.wkhtmltopdf_config is not False
        
        logger.info(f"LibreOffice available: {self.has_libreoffice}")
        logger.info(f"wkhtmltopdf available: {self.has_wkhtmltopdf}")
    
    def get_capabilities(self):
        """Return current conversion capabilities based on what's actually working"""
        capabilities = {
            'python_optimized': [],
            'libreoffice': [],
            'cloudconvert': [],
            'supported_conversions': []
        }
        
        # Python library optimized conversions
        python_conversions = [
            # Images to PDF
            'jpg_to_pdf', 'jpeg_to_pdf', 'png_to_pdf', 'gif_to_pdf', 'bmp_to_pdf', 'tiff_to_pdf',
            # Image conversions
            'png_to_jpg', 'png_to_jpeg', 'jpg_to_png', 'jpeg_to_png',
            # Images to Office/HTML/TXT
            'png_to_docx', 'jpg_to_docx', 'jpeg_to_docx',
            'png_to_pptx', 'jpg_to_pptx', 'jpeg_to_pptx',
            'png_to_html', 'jpg_to_html', 'jpeg_to_html',
            'png_to_txt', 'jpg_to_txt', 'jpeg_to_txt',
            'png_to_csv', 'jpg_to_csv', 'jpeg_to_csv',
            'png_to_xlsx', 'jpg_to_xlsx', 'jpeg_to_xlsx',
            # Office/HTML/TXT to Images (VICE-VERSA)
            'docx_to_png', 'docx_to_jpg', 'docx_to_jpeg',
            'pptx_to_png', 'pptx_to_jpg', 'pptx_to_jpeg',
            'html_to_png', 'html_to_jpg', 'html_to_jpeg',
            'csv_to_png', 'csv_to_jpg', 'csv_to_jpeg',
            'xlsx_to_png', 'xlsx_to_jpg', 'xlsx_to_jpeg',
            # PDF extractions
            'pdf_to_txt', 'pdf_to_html', 'pdf_to_csv', 'pdf_to_xlsx', 'pdf_to_xls',
            'pdf_to_jpg', 'pdf_to_jpeg', 'pdf_to_png',
            'pdf_to_docx', 'pdf_to_doc', 'pdf_to_pptx', 'pdf_to_ppt',
            # CSV conversions
            'csv_to_xlsx', 'csv_to_xls', 'csv_to_html', 'csv_to_pptx', 'csv_to_ppt',
            # TXT conversions (Python-based)
            'txt_to_html', 'txt_to_pptx', 'txt_to_csv', 'txt_to_xlsx', 'txt_to_png', 'txt_to_jpg',
            # HTML conversions (Python-based)
            'html_to_txt', 'html_to_csv', 'html_to_xlsx', 'html_to_pptx', 'html_to_png', 'html_to_jpg'
        ]
        
        if self.has_wkhtmltopdf:
            python_conversions.append('html_to_pdf')
        
        capabilities['python_optimized'] = python_conversions
        
        # LibreOffice conversions (realistic subset that actually work)
        if self.has_libreoffice:
            # Define what LibreOffice can actually convert reliably
            libreoffice_conversions = [
                # Office documents to PDF - very reliable
                'doc_to_pdf', 'docx_to_pdf', 'xls_to_pdf', 'xlsx_to_pdf', 
                'ppt_to_pdf', 'pptx_to_pdf', 'odt_to_pdf', 'ods_to_pdf', 'odp_to_pdf',
                'rtf_to_pdf', 'txt_to_pdf',
                
                # Office documents to HTML - reliable
                'doc_to_html', 'docx_to_html', 'xls_to_html', 'xlsx_to_html',
                'ppt_to_html', 'pptx_to_html', 'odt_to_html', 'ods_to_html', 'odp_to_html',
                
                # Within same document type - reliable
                'doc_to_docx', 'docx_to_doc', 'docx_to_odt', 'odt_to_docx',
                'xls_to_xlsx', 'xlsx_to_xls', 'xlsx_to_ods', 'ods_to_xlsx',
                'ppt_to_pptx', 'pptx_to_ppt', 'pptx_to_odp', 'odp_to_pptx',
                'docx_to_pptx', # Added for better reliability, LibreOffice might handle this
                
                # Text formats
                'txt_to_doc', 'txt_to_docx', 'txt_to_odt', 'rtf_to_doc', 'rtf_to_docx',
                'doc_to_rtf', 'docx_to_rtf', 'odt_to_rtf',
                
                # Some PDF input conversions (limited)
                'pdf_to_odt', 'pdf_to_rtf',
                
                # HTML conversions via LibreOffice
                'html_to_pdf', 'html_to_doc', 'html_to_docx', 'html_to_odt'
            ]
            
            capabilities['libreoffice'] = libreoffice_conversions
        
        # CloudConvert fallback
        if CLOUDCONVERT_API_KEY:
            capabilities['cloudconvert'] = ['fallback_for_unsupported_formats']
        
        # Combine all supported conversions
        all_conversions = set()
        for category in ['python_optimized', 'libreoffice']:
            all_conversions.update(capabilities[category])
        
        capabilities['supported_conversions'] = sorted(list(all_conversions))
        
        return capabilities
    
    def _cleanup_libreoffice_processes(self):
        """Kill any existing LibreOffice processes that might interfere."""
        try:
            if os.name == 'nt':  # Windows
                subprocess.run(['taskkill', '/f', '/im', 'soffice.exe'], 
                             capture_output=True, timeout=5)
                subprocess.run(['taskkill', '/f', '/im', 'soffice.bin'], 
                             capture_output=True, timeout=5)
            else:  # Linux/Mac
                subprocess.run(['pkill', '-f', 'soffice'], 
                             capture_output=True, timeout=5)
        except Exception:
            pass  # Ignore errors - this is just cleanup
    
    def convert_with_libreoffice(self, input_path, output_format, input_ext=None):
        """Convert file using LibreOffice with simplified, reliable approach."""
        if not self.has_libreoffice:
            raise Exception("LibreOffice not available")
        
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)
        
        # Validate file exists and is readable
        if not os.path.exists(input_path) or os.path.getsize(input_path) == 0:
            raise Exception(f"Input file does not exist or is empty: {input_path}")
        
        # Check if this is a supported conversion for LibreOffice
        conversion_key = f"{input_ext}_to_{output_format}"
        capabilities = self.get_capabilities()
        if conversion_key not in capabilities.get('libreoffice', []):
            raise Exception(f"LibreOffice does not support conversion from {input_ext} to {output_format}")
        
        # Clean up any existing LibreOffice processes
        self._cleanup_libreoffice_processes()
        time.sleep(0.5)  # Give processes time to fully terminate
        
        try:
            # Use absolute paths to avoid issues
            abs_input_path = os.path.abspath(input_path)
            abs_output_dir = os.path.abspath(OUTPUT_FOLDER)
            
            # First try with minimal command (most compatible)
            cmd = [
                self.libreoffice_path,
                '--headless',
                '--convert-to',
                output_format,
                '--outdir',
                abs_output_dir,
                abs_input_path
            ]
            
            # Set up environment
            env = os.environ.copy()
            # Add LibreOffice directory to PATH on Windows
            if os.name == 'nt':
                lo_dir = os.path.dirname(self.libreoffice_path)
                env['PATH'] = lo_dir + os.pathsep + env.get('PATH', '')
            
            logger.info(f"Running LibreOffice conversion: {input_ext or 'unknown'} -> {output_format}")
            logger.info(f"Command: {' '.join(cmd)}")
            
            # Run conversion with minimal setup first
            result = subprocess.run(
                cmd, 
                check=False,
                capture_output=True, 
                text=True, 
                timeout=90,
                env=env,
                cwd=abs_output_dir,  # Run from output directory
                stdin=subprocess.DEVNULL  # Prevent interactive prompts
            )
            
            # Check for success
            if result.returncode == 0:
                return True
            
            # Handle errors
            error_output = result.stderr.strip() if result.stderr else ''
            stdout_output = result.stdout.strip() if result.stdout else ''
            
            # Provide detailed error analysis
            logger.warning(f"LibreOffice stderr: {error_output}")
            logger.warning(f"LibreOffice stdout: {stdout_output}")
            
            # Check for common LibreOffice issues
            if 'source file could not be loaded' in error_output.lower():
                raise Exception(f"LibreOffice cannot read the input file format ({input_ext or 'unknown'})")
            elif 'no export filter' in error_output.lower() or 'filter' in error_output.lower():
                raise Exception(f"LibreOffice does not support conversion from {input_ext or 'unknown'} to {output_format}")
            elif 'error' in error_output.lower() and ('opening' in error_output.lower() or 'loading' in error_output.lower()):
                raise Exception(f"LibreOffice failed to open the input file: {error_output}")
            elif result.returncode == 1 and not error_output and not stdout_output:
                # Exit code 1 with no output often means unsupported conversion
                raise Exception(f"LibreOffice does not support conversion from {input_ext or 'unknown'} to {output_format} (exit code 1)")
            
            # If we get here, conversion failed
            if error_output:
                raise Exception(f"LibreOffice conversion failed: {error_output}")
            elif stdout_output:
                raise Exception(f"LibreOffice conversion failed: {stdout_output}")
            else:
                raise Exception(f"LibreOffice conversion failed with exit code {result.returncode}")
            
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out (90s)")
        finally:
            # Clean up any remaining LibreOffice processes
            self._cleanup_libreoffice_processes()
    
    def convert_image_to_pdf(self, input_path, output_path):
        """Convert image to PDF using PIL."""
        try:
            logger.info(f"Converting image {input_path} to PDF")
            image = Image.open(input_path)
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            image.save(output_path, 'PDF')
            logger.info(f"Successfully converted image to PDF: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to PDF conversion failed: {str(e)}")
            return False
    
    def convert_html_to_pdf(self, input_path, output_path):
        """Convert HTML to PDF using wkhtmltopdf."""
        if not self.has_wkhtmltopdf:
            raise Exception("wkhtmltopdf not available")
        
        try:
            logger.info(f"Converting HTML {input_path} to PDF")
            if self.wkhtmltopdf_config:
                pdfkit.from_file(input_path, output_path, configuration=self.wkhtmltopdf_config)
            else:
                pdfkit.from_file(input_path, output_path)
            
            logger.info(f"Successfully converted HTML to PDF: {output_path}")
            return True
        except Exception as e:
            logger.error(f"HTML to PDF conversion failed: {str(e)}")
            return False
    
    def convert_pdf_to_images(self, input_path, output_path, output_format='jpg'):
        """Convert PDF to images using pypdfium2 or pdf2image if available."""
        try:
            # Normalize format names
            format_map = {
                'jpg': 'JPEG',
                'jpeg': 'JPEG', 
                'png': 'PNG'
            }
            pil_format = format_map.get(output_format.lower(), output_format.upper())
            
            # Try pypdfium2 first (fast and easy on Windows)
            try:
                import pypdfium2 as pdfium
                pdf = pdfium.PdfDocument(input_path)
                page = pdf[0]
                bitmap = page.render(scale=2)
                pil_image = bitmap.to_pil()
                
                # Convert to RGB for JPEG format
                if pil_format == 'JPEG' and pil_image.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', pil_image.size, (255, 255, 255))
                    if pil_image.mode == 'P':
                        pil_image = pil_image.convert('RGBA')
                    background.paste(pil_image, mask=pil_image.split()[-1] if pil_image.mode in ('RGBA', 'LA') else None)
                    pil_image = background
                
                pil_image.save(output_path, pil_format)
                return True
            except ImportError:
                pass

            # Fallback: pdf2image if installed
            try:
                import pdf2image
                images = pdf2image.convert_from_path(input_path)
                img = images[0]
                
                # Convert to RGB for JPEG format
                if pil_format == 'JPEG' and img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = background
                    
                img.save(output_path, pil_format)
                return True
            except ImportError:
                logger.warning("Neither pypdfium2 nor pdf2image available for PDF to image conversion")
                return False
        except Exception as e:
            logger.error(f"PDF to image conversion failed: {str(e)}")
            return False

    def _pdf_to_image_fallback(self, input_path, output_path):
        """Deprecated fallback method (kept for compatibility)."""
        logger.warning("PDF to image conversion fallback not implemented")
        return False

    def convert_pdf_to_csv(self, input_path, output_path):
        """Extract tables from PDF and save as CSV using pdfplumber and pandas."""
        try:
            all_tables = []
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        if table: # Ensure table is not empty
                            df = pd.DataFrame(table[1:], columns=table[0])
                            all_tables.append(df)
            
            if all_tables:
                # Concatenate all tables into a single DataFrame if multiple tables are found
                combined_df = pd.concat(all_tables, ignore_index=True)
                combined_df.to_csv(output_path, index=False, encoding='utf-8')
                logger.info(f"Successfully converted PDF tables to CSV: {output_path}")
                return True
            else:
                logger.warning(f"No tables found in PDF {input_path} for CSV conversion.")
                return False
        except Exception as e:
            logger.error(f"PDF to CSV conversion failed: {str(e)}")
            return False

    def convert_pdf_extract_text(self, input_path, output_path, output_format):
        """Extract text or tables from PDF and convert to TXT/HTML/CSV."""
        try:
            texts = []
            tables_accum = []

            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
                    # extract tables for CSV when present
                    try:
                        page_tables = page.extract_tables()
                        for t in page_tables or []:
                            if t:
                                df = pd.DataFrame(t[1:], columns=t[0] if t[0] else None)
                                tables_accum.append(df)
                    except Exception:
                        pass

            full_text = '\n'.join(texts)
            fmt = output_format.lower()
            if fmt == 'txt':
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(full_text)
            elif fmt == 'html':
                html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset=\"utf-8\">
    <title>Extracted PDF Content</title>
    <style>body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}</style>
</head>
<body>
    <h1>Extracted PDF Content</h1>
    <pre style=\"white-space: pre-wrap; word-wrap: break-word;\">{full_text}</pre>
</body>
</html>"""
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(html_content)
            elif fmt == 'csv':
                # If we have tables, concatenate them; else, dump text lines
                import csv
                with open(output_path, 'w', encoding='utf-8', newline='') as f:
                    writer = csv.writer(f)
                    if tables_accum:
                        for i, df in enumerate(tables_accum, 1):
                            writer.writerow([f'-- Table {i} --'])
                            writer.writerow(df.columns.tolist() if list(df.columns)[0] is not None else [])
                            for _, row in df.iterrows():
                                writer.writerow(row.tolist())
                            writer.writerow([])
                    else:
                        for line in full_text.split('\n'):
                            writer.writerow([line])
            elif fmt in ['xlsx', 'xls']:
                # Create Excel file from extracted tables and text
                if tables_accum:
                    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                        for i, df in enumerate(tables_accum):
                            sheet_name = f'Table_{i+1}' if len(tables_accum) > 1 else 'Data'
                            df.to_excel(writer, sheet_name=sheet_name, index=False)
                        # Add text content as separate sheet if we have both
                        if full_text.strip():
                            text_lines = [[line] for line in full_text.split('\n') if line.strip()]
                            text_df = pd.DataFrame(text_lines, columns=['Text Content'])
                            text_df.to_excel(writer, sheet_name='Text', index=False)
                else:
                    # No tables, create simple Excel with text
                    text_lines = [[line] for line in full_text.split('\n') if line.strip()]
                    text_df = pd.DataFrame(text_lines, columns=['Text Content'])
                    text_df.to_excel(output_path, index=False)
            else:
                raise ValueError(f'Unsupported output format for text extraction: {output_format}')

            return True
        except Exception as e:
            logger.error(f"PDF text extraction failed: {str(e)}")
            return False
    
    def convert_with_cloudconvert(self, input_path, output_format):
        """Convert using CloudConvert API."""
        if not CLOUDCONVERT_API_KEY:
            raise Exception("CloudConvert API key not configured")
        
        try:
            headers = {
                'Authorization': f'Bearer {CLOUDCONVERT_API_KEY}',
                'Content-Type': 'application/json'
            }
            
            # Create job
            job_data = {
                "tasks": {
                    "import": {"operation": "import/upload"},
                    "convert": {
                        "operation": "convert",
                        "input": "import",
                        "output_format": output_format
                    },
                    "export": {"operation": "export/url", "input": "convert"}
                }
            }
            
            response = requests.post('https://api.cloudconvert.com/v2/jobs', 
                                   headers=headers, json=job_data, timeout=30)
            response.raise_for_status()
            
            job = response.json()
            job_id = job['data']['id']
            import_task = job['data']['tasks'][0]
            
            # Upload file
            with open(input_path, 'rb') as f:
                files = {'file': f}
                upload_response = requests.post(
                    import_task['result']['form']['url'],
                    data=import_task['result']['form']['parameters'],
                    files=files, timeout=120
                )
                upload_response.raise_for_status()
            
            # Wait for completion and download
            max_wait = 300  # 5 minutes
            wait_time = 0
            while wait_time < max_wait:
                status_response = requests.get(f'https://api.cloudconvert.com/v2/jobs/{job_id}', 
                                             headers=headers, timeout=30)
                status_response.raise_for_status()
                
                job_status = status_response.json()
                if job_status['data']['status'] == 'finished':
                    export_task = next(t for t in job_status['data']['tasks'] if t['name'] == 'export')
                    download_url = export_task['result']['files'][0]['url']
                    
                    download_response = requests.get(download_url, timeout=120)
                    download_response.raise_for_status()
                    
                    return download_response.content
                elif job_status['data']['status'] == 'error':
                    raise Exception("CloudConvert conversion failed")
                
                time.sleep(5)
                wait_time += 5
            
            raise Exception("CloudConvert conversion timed out")
            
        except Exception as e:
            logger.error(f"CloudConvert conversion failed: {str(e)}")
            raise
    
    def convert_pdf_to_word(self, input_path, output_path, output_format='docx'):
        """Convert PDF to Word document by extracting text and creating document."""
        try:
            # Check if python-docx is available
            try:
                from docx import Document
                from docx.shared import Inches
            except ImportError:
                # Fallback to LibreOffice if python-docx not available
                logger.warning("python-docx not available, trying LibreOffice for PDF to Word conversion")
                return False
            
            logger.info(f"Converting PDF {input_path} to {output_format.upper()} using text extraction")
            
            # Extract text from PDF
            all_text = []
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        all_text.append(f"--- Page {page_num + 1} ---\n{text.strip()}\n")
            
            if not all_text:
                logger.warning("No text found in PDF for conversion")
                return False
            
            # Create Word document
            doc = Document()
            doc.add_heading('PDF Content', 0)
            
            full_text = '\n'.join(all_text)
            
            # Split text into paragraphs and add to document
            paragraphs = full_text.split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    if paragraph.strip().startswith('--- Page'):
                        # Add page headings
                        doc.add_heading(paragraph.strip(), level=1)
                    else:
                        doc.add_paragraph(paragraph.strip())
            
            # Save the document
            doc.save(output_path)
            logger.info(f"Successfully converted PDF to {output_format.upper()}: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PDF to {output_format.upper()} conversion failed: {str(e)}")
            return False
    
    def convert_pdf_to_pptx(self, input_path, output_path):
        """Convert PDF to PPTX by extracting text and creating slides."""
        try:
            # Check if python-pptx is available
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
            except ImportError:
                logger.error("python-pptx not available for PDF to PowerPoint conversion")
                return False
            
            logger.info(f"Converting PDF {input_path} to PPTX using text extraction")
            
            # Extract text from PDF
            texts = []
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        texts.append((page_num + 1, text.strip()))
            
            if not texts:
                logger.warning("No text found in PDF for conversion")
                return False
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = "PDF Content"
            subtitle.text = f"Converted from PDF ({len(texts)} pages)"
            
            # Add content slides
            for page_num, page_text in texts:
                # Use content slide layout
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Set slide title
                title = slide.shapes.title
                title.text = f"Page {page_num}"
                
                # Add text content
                content = slide.placeholders[1]
                
                # Split text into manageable chunks (PowerPoint has character limits)
                max_chars_per_slide = 1000
                if len(page_text) > max_chars_per_slide:
                    # Split long text into multiple parts
                    words = page_text.split()
                    chunks = []
                    current_chunk = []
                    current_length = 0
                    
                    for word in words:
                        if current_length + len(word) + 1 > max_chars_per_slide and current_chunk:
                            chunks.append(' '.join(current_chunk))
                            current_chunk = [word]
                            current_length = len(word)
                        else:
                            current_chunk.append(word)
                            current_length += len(word) + 1
                    
                    if current_chunk:
                        chunks.append(' '.join(current_chunk))
                    
                    # Use first chunk for this slide
                    content.text = chunks[0]
                    
                    # Create additional slides for remaining chunks
                    for i, chunk in enumerate(chunks[1:], 1):
                        extra_slide = prs.slides.add_slide(prs.slide_layouts[1])
                        extra_title = extra_slide.shapes.title
                        extra_title.text = f"Page {page_num} (Part {i+1})"
                        extra_content = extra_slide.placeholders[1]
                        extra_content.text = chunk
                else:
                    content.text = page_text
            
            # Save the presentation
            prs.save(output_path)
            logger.info(f"Successfully converted PDF to PPTX: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PDF to PPTX conversion failed: {str(e)}")
            return False
    
    def convert_pptx_to_pdf(self, input_path, output_path):
        """Convert PPTX to PDF by extracting text and creating PDF document."""
        try:
            # Check if required libraries are available
            try:
                from pptx import Presentation
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import inch
            except ImportError:
                # Fallback to basic PDF creation with just text
                try:
                    from pptx import Presentation
                    # Use the existing convert_image_to_pdf as a template for basic PDF creation
                    return self._convert_pptx_to_pdf_basic(input_path, output_path)
                except ImportError as e:
                    logger.warning(f"Required libraries not available for PPTX to PDF conversion: {e}")
                    return False
            
            logger.info(f"Converting PPTX {input_path} to PDF using text extraction")
            
            # Load PowerPoint presentation
            prs = Presentation(input_path)
            
            # Create PDF document
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
            )
            story.append(Paragraph("PowerPoint Content", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            slide_count = 0
            for i, slide in enumerate(prs.slides):
                slide_count += 1
                slide_title = f"Slide {slide_count}"
                
                # Try to get slide title from title shape
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if len(shape.text.strip()) < 100:  # Likely a title
                            slide_title = shape.text.strip()
                            break
                
                # Add slide heading
                story.append(Paragraph(slide_title, styles['Heading2']))
                story.append(Spacer(1, 0.1*inch))
                
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if text not in slide_text and text != slide_title:
                            slide_text.append(text)
                
                # Add slide content
                if slide_text:
                    for text in slide_text:
                        story.append(Paragraph(text, styles['Normal']))
                        story.append(Spacer(1, 0.1*inch))
                else:
                    story.append(Paragraph("[No text content found in this slide]", styles['Italic']))
                
                # Add spacing between slides
                story.append(Spacer(1, 0.3*inch))
            
            if slide_count == 0:
                story.append(Paragraph("[No slides found in the PowerPoint presentation]", styles['Italic']))
            
            # Build PDF
            doc.build(story)
            logger.info(f"Successfully converted PPTX to PDF: {output_path} ({slide_count} slides)")
            return True
            
        except Exception as e:
            logger.error(f"PPTX to PDF conversion failed: {str(e)}")
            return False
    
    def _convert_pptx_to_pdf_basic(self, input_path, output_path):
        """Basic PPTX to PDF conversion fallback using PIL for simple text rendering."""
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw, ImageFont
            import io
            
            # Load PowerPoint presentation
            prs = Presentation(input_path)
            
            # Extract all text
            all_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    all_text.extend(slide_text)
            
            if not all_text:
                all_text = ["[No text found in PowerPoint presentation]"]
            
            # Create a simple text-based image and convert to PDF
            img = Image.new('RGB', (2480, 3508), color='white')  # A4 size at 300 DPI
            draw = ImageDraw.Draw(img)
            
            try:
                font = ImageFont.truetype("arial.ttf", 40)
            except:
                font = ImageFont.load_default()
            
            y_position = 100
            for text in all_text:
                # Word wrap
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    current_line.append(word)
                    line_text = ' '.join(current_line)
                    if draw.textlength(line_text, font=font) > 2300:  # Leave margin
                        if len(current_line) > 1:
                            lines.append(' '.join(current_line[:-1]))
                            current_line = [word]
                        else:
                            lines.append(line_text)
                            current_line = []
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines:
                    if y_position > 3300:  # Near bottom of page
                        break
                    draw.text((100, y_position), line, font=font, fill='black')
                    y_position += 60
                
                y_position += 40  # Extra space between text blocks
            
            # Convert image to PDF
            img.save(output_path, 'PDF')
            logger.info(f"Successfully converted PPTX to PDF (basic method): {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Basic PPTX to PDF conversion failed: {str(e)}")
            return False
    
    def convert_excel_to_word(self, input_path, output_path, output_format='docx'):
        """Convert Excel to Word by extracting data and creating document with tables."""
        try:
            # Check if required libraries are available
            try:
                from docx import Document
                from docx.shared import Inches
            except ImportError as e:
                logger.warning(f"Required libraries not available for Excel to Word conversion: {e}")
                return False
            
            logger.info(f"Converting Excel {input_path} to {output_format.upper()} using data extraction")
            
            # Read Excel file using pandas
            try:
                # Try to read all sheets
                excel_file = pd.ExcelFile(input_path)
                sheet_names = excel_file.sheet_names
            except Exception as e:
                logger.error(f"Failed to read Excel file: {e}")
                return False
            
            # Create Word document
            doc = Document()
            doc.add_heading('Excel Content', 0)
            
            sheet_count = 0
            for sheet_name in sheet_names:
                try:
                    # Read the sheet
                    df = pd.read_excel(input_path, sheet_name=sheet_name)
                    
                    if df.empty:
                        continue
                    
                    sheet_count += 1
                    
                    # Add sheet heading
                    if len(sheet_names) > 1:
                        doc.add_heading(f'Sheet: {sheet_name}', level=1)
                    
                    # Limit the size of tables for Word compatibility
                    max_rows = 50
                    max_cols = 10
                    
                    if len(df) > max_rows:
                        doc.add_paragraph(f"Note: Showing first {max_rows} rows of {len(df)} total rows.")
                        df = df.head(max_rows)
                    
                    if len(df.columns) > max_cols:
                        doc.add_paragraph(f"Note: Showing first {max_cols} columns of {len(df.columns)} total columns.")
                        df = df.iloc[:, :max_cols]
                    
                    # Create table in Word document
                    if not df.empty:
                        rows = len(df) + 1  # +1 for header
                        cols = len(df.columns)
                        table = doc.add_table(rows=rows, cols=cols)
                        table.style = 'Table Grid'
                        
                        # Add headers
                        for j, column_name in enumerate(df.columns):
                            cell = table.cell(0, j)
                            cell.text = str(column_name)
                            # Make header bold
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                        
                        # Add data rows
                        for i, (_, row) in enumerate(df.iterrows()):
                            for j, value in enumerate(row):
                                table.cell(i + 1, j).text = str(value) if pd.notna(value) else ""
                    
                    # Add some spacing between sheets
                    doc.add_paragraph()
                    
                except Exception as e:
                    logger.warning(f"Failed to process sheet '{sheet_name}': {e}")
                    doc.add_paragraph(f"[Error processing sheet '{sheet_name}': {e}]")
                    continue
            
            if sheet_count == 0:
                logger.warning("No valid sheets found in Excel file")
                doc.add_paragraph("[No valid data sheets found in the Excel file]")
            
            # Save the document
            doc.save(output_path)
            logger.info(f"Successfully converted Excel to {output_format.upper()}: {output_path} ({sheet_count} sheets)")
            return True
            
        except Exception as e:
            logger.error(f"Excel to {output_format.upper()} conversion failed: {str(e)}")
            return False
    
    def convert_docx_to_pptx(self, input_path, output_path):
        """Convert DOCX to PPTX by extracting text and creating slides."""
        try:
            # Check if required libraries are available
            try:
                from docx import Document
                from pptx import Presentation
                from pptx.util import Inches, Pt
            except ImportError as e:
                logger.warning(f"Required libraries not available for DOCX to PPTX conversion: {e}")
                return False
            
            logger.info(f"Converting DOCX {input_path} to PPTX using text extraction")
            
            # Try to load Word document with better error handling
            try:
                doc = Document(input_path)
            except Exception as e:
                error_msg = str(e).lower()
                logger.warning(f"Failed to open DOCX file: {error_msg}")
                
                if 'spreadsheet' in error_msg or 'xlsx' in error_msg or 'sheet' in error_msg:
                    raise Exception("File appears to be an Excel spreadsheet, not a Word document. Please verify the file type and rename if necessary.")
                elif 'presentation' in error_msg or 'pptx' in error_msg:
                    raise Exception("File appears to be a PowerPoint presentation, not a Word document. Please verify the file type.")
                elif 'not a word file' in error_msg:
                    raise Exception("File is not a valid Word document. It may be corrupted or have the wrong file extension.")
                else:
                    raise Exception(f"Cannot open file as Word document. The file may be corrupted, password-protected, or not a valid DOCX file. Error: {str(e)}")
            
            # Extract text from document
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            if not paragraphs:
                logger.warning("No text found in Word document")
                # Create a placeholder presentation for empty document
                prs = Presentation()
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_slide.shapes.title.text = "Empty Document"
                title_slide.placeholders[1].text = "No text content found in the Word document"
                
                # Add a content slide explaining the situation
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "No Content"
                slide.placeholders[1].text = "The Word document appears to be empty or contains only formatting without readable text."
                
                prs.save(output_path)
                logger.info(f"Created empty PPTX for empty DOCX: {output_path}")
                return True
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            # Use first paragraph as title, or default
            if paragraphs:
                first_para = paragraphs[0]
                if len(first_para) < 100:  # Likely a title
                    title.text = first_para
                    subtitle.text = "Converted from Word Document"
                    content_paragraphs = paragraphs[1:]
                else:
                    title.text = "Word Document Content"
                    subtitle.text = "Converted from DOCX"
                    content_paragraphs = paragraphs
            else:
                title.text = "Word Document Content"
                subtitle.text = "Converted from DOCX"
                content_paragraphs = []
            
            # Add content slides
            max_chars_per_slide = 800
            current_slide_content = []
            current_slide_chars = 0
            
            for paragraph in content_paragraphs:
                # Check if adding this paragraph would exceed the character limit
                if current_slide_chars + len(paragraph) > max_chars_per_slide and current_slide_content:
                    # Create slide with current content
                    self._create_content_slide(prs, current_slide_content)
                    current_slide_content = [paragraph]
                    current_slide_chars = len(paragraph)
                else:
                    current_slide_content.append(paragraph)
                    current_slide_chars += len(paragraph) + 1  # +1 for newline
            
            # Create final slide if there's remaining content
            if current_slide_content:
                self._create_content_slide(prs, current_slide_content)
            
            # Ensure we have at least one content slide
            if len(prs.slides) == 1:  # Only title slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Content"
                slide.placeholders[1].text = "[No content paragraphs found in the Word document]"
            
            # Save the presentation
            prs.save(output_path)
            logger.info(f"Successfully converted DOCX to PPTX: {output_path} ({len(prs.slides)} slides)")
            return True
            
        except Exception as e:
            logger.error(f"DOCX to PPTX conversion failed: {str(e)}")
            return False
    
    def _create_content_slide(self, prs, content_list):
        """Helper method to create a content slide with text."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide layout
        
        # Set title based on first line if it's short, otherwise generic title
        first_line = content_list[0] if content_list else "Content"
        if len(first_line) < 60:
            slide.shapes.title.text = first_line
            slide_content = content_list[1:] if len(content_list) > 1 else []
        else:
            slide.shapes.title.text = "Content"
            slide_content = content_list
        
        # Add content
        if slide_content:
            content_text = '\n\n'.join(slide_content)
        else:
            content_text = first_line if len(first_line) >= 60 else "[Content slide]"
        
        slide.placeholders[1].text = content_text
    
    def convert_pptx_to_word(self, input_path, output_path, output_format='docx'):
        """Convert PPTX to Word by extracting text and creating document."""
        try:
            # Check if required libraries are available
            try:
                from pptx import Presentation
                from docx import Document
                from docx.shared import Inches
            except ImportError as e:
                logger.warning(f"Required libraries not available for PPTX to Word conversion: {e}")
                return False
            
            logger.info(f"Converting PPTX {input_path} to {output_format.upper()} using text extraction")
            
            # Load PowerPoint presentation
            prs = Presentation(input_path)
            
            # Create Word document
            doc = Document()
            doc.add_heading('PowerPoint Content', 0)
            
            slide_count = 0
            for i, slide in enumerate(prs.slides):
                slide_count += 1
                slide_title = f"Slide {slide_count}"
                
                # Try to get slide title from title shape
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if len(shape.text.strip()) < 100:  # Likely a title
                            slide_title = shape.text.strip()
                            break
                
                # Add slide heading
                doc.add_heading(slide_title, level=1)
                
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if text not in slide_text:  # Avoid duplicates
                            slide_text.append(text)
                
                # Add slide content
                if slide_text:
                    for text in slide_text:
                        if text != slide_title:  # Don't repeat the title
                            doc.add_paragraph(text)
                else:
                    doc.add_paragraph("[No text content found in this slide]")
                
                # Add some spacing between slides
                doc.add_paragraph()
            
            if slide_count == 0:
                logger.warning("No slides found in PowerPoint presentation")
                doc.add_paragraph("[No slides found in the PowerPoint presentation]")
            
            # Save the document
            doc.save(output_path)
            logger.info(f"Successfully converted PPTX to {output_format.upper()}: {output_path} ({slide_count} slides)")
            return True
            
        except Exception as e:
            logger.error(f"PPTX to {output_format.upper()} conversion failed: {str(e)}")
            return False
    
    def convert_txt_to_html(self, input_path, output_path):
        """Convert TXT to HTML with proper formatting."""
        try:
            # Read the text file
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Simple HTML conversion with basic formatting
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Text Document</title>
    <style>
        body {{
            font-family: 'Courier New', monospace;
            line-height: 1.6;
            margin: 40px;
            background-color: #f9f9f9;
        }}
        .content {{
            background-color: white;
            padding: 30px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            white-space: pre-wrap;
            word-wrap: break-word;
        }}
        h1 {{
            color: #333;
            border-bottom: 2px solid #ccc;
            padding-bottom: 10px;
        }}
    </style>
</head>
<body>
    <h1>Text Document</h1>
    <div class="content">{content}</div>
</body>
</html>"""
            
            # Write HTML file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info(f"Successfully converted TXT to HTML: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"TXT to HTML conversion failed: {str(e)}")
            return False
    
    def convert_txt_to_pptx(self, input_path, output_path):
        """Convert TXT to PPTX by creating slides from text content."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Read the text file
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            if not content.strip():
                logger.warning("Empty text file")
                return False
            
            # Split content into paragraphs
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            if not paragraphs:
                paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Text Document"
            title_slide.placeholders[1].text = "Converted from TXT file"
            
            # Content slides
            max_chars_per_slide = 800
            current_content = []
            current_chars = 0
            
            for para in paragraphs:
                if current_chars + len(para) > max_chars_per_slide and current_content:
                    self._create_txt_content_slide(prs, current_content)
                    current_content = [para]
                    current_chars = len(para)
                else:
                    current_content.append(para)
                    current_chars += len(para) + 2  # +2 for paragraph breaks
            
            if current_content:
                self._create_txt_content_slide(prs, current_content)
            
            # Ensure at least one content slide
            if len(prs.slides) == 1:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Content"
                slide.placeholders[1].text = content[:1000] + ("..." if len(content) > 1000 else "")
            
            prs.save(output_path)
            logger.info(f"Successfully converted TXT to PPTX: {output_path} ({len(prs.slides)} slides)")
            return True
            
        except Exception as e:
            logger.error(f"TXT to PPTX conversion failed: {str(e)}")
            return False
    
    def _create_txt_content_slide(self, prs, content_list):
        """Helper to create content slide from text paragraphs."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Use first paragraph as title if it's short
        first_para = content_list[0] if content_list else "Content"
        if len(first_para) < 50 and '=' not in first_para and '-' not in first_para:
            slide.shapes.title.text = first_para
            content = content_list[1:] if len(content_list) > 1 else []
        else:
            slide.shapes.title.text = "Content"
            content = content_list
        
        if content:
            slide.placeholders[1].text = '\n\n'.join(content)
        else:
            slide.placeholders[1].text = first_para
    
    def convert_html_to_txt(self, input_path, output_path):
        """Convert HTML to TXT by extracting text content."""
        try:
            # Try to use BeautifulSoup for better HTML parsing
            try:
                from bs4 import BeautifulSoup
                with open(input_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                text = soup.get_text(separator='\n\n')
                
            except ImportError:
                # Fallback: simple regex-based HTML tag removal
                import re
                with open(input_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # Remove script and style elements
                html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
                html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
                
                # Remove HTML tags
                text = re.sub(r'<[^>]+>', '', html_content)
                
                # Clean up whitespace
                text = re.sub(r'\n\s*\n', '\n\n', text)
                text = re.sub(r'[ \t]+', ' ', text)
            
            # Clean up the text (import re for this section)
            import re
            text = '\n'.join(line.strip() for line in text.split('\n'))
            text = re.sub(r'\n{3,}', '\n\n', text)
            text = text.strip()
            
            # Write text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            logger.info(f"Successfully converted HTML to TXT: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"HTML to TXT conversion failed: {str(e)}")
            return False
    
    def convert_html_to_pptx(self, input_path, output_path):
        """Convert HTML to PPTX by extracting content and creating slides."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # First convert HTML to text
            try:
                from bs4 import BeautifulSoup
                with open(input_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract title
                title_elem = soup.find('title')
                title = title_elem.get_text().strip() if title_elem else "HTML Document"
                
                # Extract headings and content
                sections = []
                
                # Get all headings and paragraphs in order
                for elem in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'li']):
                    text = elem.get_text().strip()
                    if text:
                        tag_name = elem.name
                        if tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                            sections.append(('heading', text))
                        else:
                            sections.append(('content', text))
                
            except ImportError:
                # Fallback without BeautifulSoup
                import re
                with open(input_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # Extract title
                title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
                title = title_match.group(1).strip() if title_match else "HTML Document"
                
                # Simple extraction
                text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
                text = re.sub(r'<[^>]+>', '\n', text)
                text = re.sub(r'\n\s*\n', '\n', text)
                
                paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
                sections = [('content', p) for p in paragraphs]
            
            if not sections:
                logger.warning("No content found in HTML file")
                return False
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            title_slide.placeholders[1].text = "Converted from HTML"
            
            # Process sections into slides
            current_slide_content = []
            current_slide_title = "Content"
            max_chars = 800
            current_chars = 0
            
            for section_type, text in sections:
                if section_type == 'heading' and len(text) < 100:
                    # Start new slide with this heading
                    if current_slide_content:
                        self._create_html_content_slide(prs, current_slide_title, current_slide_content)
                    current_slide_title = text
                    current_slide_content = []
                    current_chars = 0
                else:
                    # Add content to current slide
                    if current_chars + len(text) > max_chars and current_slide_content:
                        self._create_html_content_slide(prs, current_slide_title, current_slide_content)
                        current_slide_title = "Content" if section_type != 'heading' else text
                        current_slide_content = [] if section_type == 'heading' else [text]
                        current_chars = 0 if section_type == 'heading' else len(text)
                    else:
                        current_slide_content.append(text)
                        current_chars += len(text) + 2
            
            # Add final slide if needed
            if current_slide_content:
                self._create_html_content_slide(prs, current_slide_title, current_slide_content)
            
            # Ensure at least one content slide
            if len(prs.slides) == 1:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Content"
                slide.placeholders[1].text = "[No content extracted from HTML]"
            
            prs.save(output_path)
            logger.info(f"Successfully converted HTML to PPTX: {output_path} ({len(prs.slides)} slides)")
            return True
            
        except Exception as e:
            logger.error(f"HTML to PPTX conversion failed: {str(e)}")
            return False
    
    def _create_html_content_slide(self, prs, title, content_list):
        """Helper to create content slide from HTML sections."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        if content_list:
            slide.placeholders[1].text = '\n\n'.join(content_list)
        else:
            slide.placeholders[1].text = "[No content]"
    
    def convert_image_formats(self, input_path, output_path, input_format, output_format):
        """Convert between image formats using PIL."""
        try:
            from PIL import Image
            
            # Open the image
            with Image.open(input_path) as img:
                # Convert mode if necessary
                if output_format.upper() == 'JPEG' and img.mode in ('RGBA', 'LA', 'P'):
                    # Convert to RGB for JPEG (no alpha channel)
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = background
                elif output_format.upper() == 'PNG' and img.mode not in ('RGBA', 'RGB', 'L', 'LA'):
                    img = img.convert('RGBA')
                
                # Save with appropriate options
                if output_format.upper() == 'JPEG':
                    img.save(output_path, 'JPEG', quality=85, optimize=True)
                else:
                    img.save(output_path, output_format.upper())
            
            logger.info(f"Successfully converted {input_format.upper()} to {output_format.upper()}: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Image conversion failed: {str(e)}")
            return False
    
    def convert_image_to_docx(self, input_path, output_path):
        """Embed image into a DOCX file, scaled to fit the page."""
        try:
            from docx import Document
            from docx.shared import Inches
            from PIL import Image as PILImage
            
            # Create document
            doc = Document()
            section = doc.sections[0]
            page_width = section.page_width.inches
            page_height = section.page_height.inches
            left_margin = section.left_margin.inches
            right_margin = section.right_margin.inches
            top_margin = section.top_margin.inches
            bottom_margin = section.bottom_margin.inches
            
            max_width = page_width - left_margin - right_margin
            max_height = page_height - top_margin - bottom_margin
            
            # Determine image size and scale
            with PILImage.open(input_path) as img:
                w, h = img.size
                dpi_scale = 96.0  # Approximate pixels per inch for screen images
                img_width_in = w / dpi_scale
                img_height_in = h / dpi_scale
                scale = min(max_width / img_width_in, max_height / img_height_in, 1.0)
                display_width = img_width_in * scale
            
            # Add image centered
            p = doc.add_paragraph()
            run = p.add_run()
            run.add_picture(input_path, width=Inches(display_width))
            
            doc.save(output_path)
            logger.info(f"Successfully converted image to DOCX: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to DOCX conversion failed: {str(e)}")
            return False
    
    def convert_image_to_pptx(self, input_path, output_path):
        """Embed image into a PPTX presentation on a single slide."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from PIL import Image as PILImage
            
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
            
            slide_width = prs.slide_width / 914400.0  # EMU to inches
            slide_height = prs.slide_height / 914400.0
            
            with PILImage.open(input_path) as img:
                w, h = img.size
                dpi_scale = 96.0
                img_w_in = w / dpi_scale
                img_h_in = h / dpi_scale
                scale = min(slide_width / img_w_in, slide_height / img_h_in, 1.0)
                disp_w = img_w_in * scale
                disp_h = img_h_in * scale
                left = (slide_width - disp_w) / 2.0
                top = (slide_height - disp_h) / 2.0
            
            slide.shapes.add_picture(input_path, Inches(left), Inches(top), width=Inches(disp_w))
            prs.save(output_path)
            logger.info(f"Successfully converted image to PPTX: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to PPTX conversion failed: {str(e)}")
            return False
    
    def convert_image_to_html(self, input_path, output_path):
        """Create an HTML file embedding the image as base64."""
        try:
            import base64
            with open(input_path, 'rb') as f:
                b64 = base64.b64encode(f.read()).decode('ascii')
            ext = os.path.splitext(input_path)[1].lstrip('.').lower()
            html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset=\"utf-8\" />
<title>Image Document</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 40px; }}
img {{ max-width: 100%; height: auto; border: 1px solid #ddd; }}
</style>
</head>
<body>
<h1>Embedded Image</h1>
<img src=\"data:image/{ext};base64,{b64}\" alt=\"Embedded Image\" />
</body>
</html>"""
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html)
            logger.info(f"Successfully converted image to HTML: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to HTML conversion failed: {str(e)}")
            return False
    
    def convert_image_to_txt(self, input_path, output_path):
        """Extract text from image using OCR if available; else write metadata."""
        try:
            text = None
            try:
                import pytesseract
                from PIL import Image as PILImage
                with PILImage.open(input_path) as img:
                    text = pytesseract.image_to_string(img)
            except Exception as ocr_err:
                logger.warning(f"OCR not available or failed: {ocr_err}")
            
            if text and text.strip():
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text.strip())
                logger.info(f"Successfully extracted text from image to TXT: {output_path}")
                return True
            else:
                # Fallback: write metadata
                from PIL import Image as PILImage
                with PILImage.open(input_path) as img:
                    w, h = img.size
                    info = img.info
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(f"[OCR unavailable] Image size: {w}x{h} pixels. Metadata: {info}\n")
                logger.info(f"Wrote image metadata to TXT: {output_path}")
                return True
        except Exception as e:
            logger.error(f"Image to TXT conversion failed: {str(e)}")
            return False
    
    def convert_txt_to_image(self, input_path, output_path, output_format='png'):
        """Render text content to an image using PIL."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            # Settings
            width = 1200
            margin = 50
            bg_color = 'white'
            text_color = 'black'
            
            try:
                font = ImageFont.truetype('arial.ttf', 20)
            except:
                font = ImageFont.load_default()
            
            # Word-wrap
            def wrap_text(draw, text, font, max_width):
                words = text.split()
                lines = []
                line = ''
                for w in words:
                    test = (line + ' ' + w).strip()
                    if draw.textlength(test, font=font) <= max_width:
                        line = test
                    else:
                        if line:
                            lines.append(line)
                        line = w
                if line:
                    lines.append(line)
                return lines
            
            tmp_img = Image.new('RGB', (width, 1))
            draw = ImageDraw.Draw(tmp_img)
            lines = []
            for paragraph in text.split('\n'):
                if paragraph.strip():
                    lines.extend(wrap_text(draw, paragraph, font, width - 2*margin))
                else:
                    lines.append('')
            line_height = int(font.size * 1.6) if hasattr(font, 'size') else 24
            height = margin*2 + line_height * max(1, len(lines))
            img = Image.new('RGB', (width, height), color=bg_color)
            d = ImageDraw.Draw(img)
            y = margin
            for line in lines:
                d.text((margin, y), line, font=font, fill=text_color)
                y += line_height
            
            fmt = 'JPEG' if output_format.lower() in ['jpg','jpeg'] else 'PNG'
            img.save(output_path, fmt)
            logger.info(f"Successfully converted TXT to {fmt}: {output_path}")
            return True
        except Exception as e:
            logger.error(f"TXT to Image conversion failed: {str(e)}")
            return False
    
    def convert_html_to_image(self, input_path, output_path, output_format='png'):
        """Simplified HTML to image: extracts text and renders as image (no CSS)."""
        try:
            # Extract text
            tmp_txt = os.path.join(tempfile.gettempdir(), f"html_text_{uuid.uuid4()}.txt")
            if self.convert_html_to_txt(input_path, tmp_txt):
                ok = self.convert_txt_to_image(tmp_txt, output_path, output_format)
                try:
                    os.remove(tmp_txt)
                except Exception:
                    pass
                return ok
            return False
        except Exception as e:
            logger.error(f"HTML to Image conversion failed: {str(e)}")
            return False
    
    def convert_image_to_csv(self, input_path, output_path):
        """Convert image to CSV by extracting tabular data or creating metadata table."""
        try:
            from PIL import Image as PILImage
            import csv
            
            # Try OCR first for actual table extraction
            text_data = None
            try:
                import pytesseract
                with PILImage.open(input_path) as img:
                    text_data = pytesseract.image_to_string(img)
                
                # Simple table detection - look for rows with consistent separators
                lines = [line.strip() for line in text_data.split('\n') if line.strip()]
                table_rows = []
                
                for line in lines:
                    # Try to split by common separators
                    if '\t' in line:
                        table_rows.append(line.split('\t'))
                    elif '  ' in line:  # Multiple spaces
                        table_rows.append([col.strip() for col in line.split('  ') if col.strip()])
                    elif '|' in line:
                        table_rows.append([col.strip() for col in line.split('|') if col.strip()])
                    else:
                        # Single row data
                        table_rows.append([line])
                
                if table_rows and len(table_rows) > 1:
                    with open(output_path, 'w', newline='', encoding='utf-8') as f:
                        writer = csv.writer(f)
                        for row in table_rows:
                            writer.writerow(row)
                    logger.info(f"Successfully extracted table data from image to CSV: {output_path}")
                    return True
            except Exception as ocr_err:
                logger.warning(f"OCR table extraction failed: {ocr_err}")
            
            # Fallback: Create metadata CSV
            with PILImage.open(input_path) as img:
                w, h = img.size
                format_name = img.format or 'Unknown'
                mode = img.mode
                info = img.info
            
            with open(output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerow(['Property', 'Value'])
                writer.writerow(['Width', w])
                writer.writerow(['Height', h])
                writer.writerow(['Format', format_name])
                writer.writerow(['Mode', mode])
                writer.writerow(['File Size', f"{os.path.getsize(input_path)} bytes"])
                for key, value in info.items():
                    writer.writerow([key, str(value)])
            
            logger.info(f"Created image metadata CSV: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to CSV conversion failed: {str(e)}")
            return False
    
    def convert_image_to_xlsx(self, input_path, output_path):
        """Convert image to XLSX by creating metadata spreadsheet."""
        try:
            from PIL import Image as PILImage
            import pandas as pd
            
            with PILImage.open(input_path) as img:
                w, h = img.size
                format_name = img.format or 'Unknown'
                mode = img.mode
                info = img.info
            
            # Create metadata DataFrame
            metadata = [
                ['Width', w],
                ['Height', h],
                ['Format', format_name],
                ['Mode', mode],
                ['File Size', f"{os.path.getsize(input_path)} bytes"]
            ]
            
            for key, value in info.items():
                metadata.append([key, str(value)])
            
            df = pd.DataFrame(metadata, columns=['Property', 'Value'])
            df.to_excel(output_path, index=False, sheet_name='Image Metadata')
            
            logger.info(f"Created image metadata XLSX: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Image to XLSX conversion failed: {str(e)}")
            return False
    
    def convert_docx_to_image(self, input_path, output_path, output_format='png'):
        """Convert DOCX to image by extracting text and rendering as image."""
        try:
            from docx import Document
            
            # Extract all text from document
            doc = Document(input_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                # Skip placeholder/empty content indicators
                if text and not any(skip in text.lower() for skip in [
                    '[empty', 'empty document', 'no content', 'no text content',
                    'powerpoint content', 'converted from', '[no content'
                ]):
                    text_content.append(text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    row_texts = [t for t in row_texts if t and not any(skip in t.lower() for skip in [
                        '[empty', 'empty document', 'no content'
                    ])]
                    if row_texts:
                        text_content.append(' | '.join(row_texts))
            
            # If still no meaningful content, try to extract from the original source
            if not text_content:
                # Get the original filename to provide better context
                filename = os.path.basename(input_path)
                text_content = [f'Document: {filename}', 'This document appears to contain formatting or images but no extractable text content.']
            
            # Create temporary text file and render as image
            import tempfile
            temp_txt = os.path.join(tempfile.gettempdir(), f"docx_text_{uuid.uuid4()}.txt")
            try:
                with open(temp_txt, 'w', encoding='utf-8') as f:
                    f.write('\n\n'.join(text_content))
                
                result = self.convert_txt_to_image(temp_txt, output_path, output_format)
                return result
            finally:
                try:
                    os.remove(temp_txt)
                except Exception:
                    pass
                    
        except Exception as e:
            logger.error(f"DOCX to image conversion failed: {str(e)}")
            return False
    
    def convert_pptx_to_image(self, input_path, output_path, output_format='png'):
        """Convert PPTX to image by extracting text and rendering as image."""
        try:
            from pptx import Presentation
            
            # Extract all text from presentation
            prs = Presentation(input_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                meaningful_content = False
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        # Skip placeholder/empty content indicators
                        if not any(skip in text.lower() for skip in [
                            '[empty', 'empty document', 'no content', 'no text content',
                            'converted from', '[no content', 'slide 1', 'slide 2'
                        ]):
                            slide_text.append(text)
                            meaningful_content = True
                
                if meaningful_content:
                    if slide_text:
                        text_content.append(f"=== SLIDE {i} ===")
                        text_content.extend(slide_text)
                        text_content.append('')  # Empty line between slides
            
            if not text_content:
                filename = os.path.basename(input_path)
                text_content = [f'Presentation: {filename}', 'This PowerPoint contains slides with formatting or images but no extractable text content.']
            
            # Create temporary text file and render as image
            import tempfile
            temp_txt = os.path.join(tempfile.gettempdir(), f"pptx_text_{uuid.uuid4()}.txt")
            try:
                with open(temp_txt, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(text_content))
                
                result = self.convert_txt_to_image(temp_txt, output_path, output_format)
                return result
            finally:
                try:
                    os.remove(temp_txt)
                except Exception:
                    pass
                    
        except Exception as e:
            logger.error(f"PPTX to image conversion failed: {str(e)}")
            return False
    
    def convert_csv_to_image(self, input_path, output_path, output_format='png'):
        """Convert CSV to image by rendering table data as image."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            import pandas as pd
            
            # Read CSV data
            df = pd.read_csv(input_path)
            
            # Settings
            cell_padding = 10
            row_height = 30
            max_col_width = 150
            bg_color = 'white'
            text_color = 'black'
            border_color = 'gray'
            header_bg = 'lightblue'
            
            try:
                font = ImageFont.truetype('arial.ttf', 12)
            except:
                font = ImageFont.load_default()
            
            # Calculate dimensions
            cols = len(df.columns)
            rows = len(df) + 1  # +1 for header
            
            # Truncate data if too large
            if rows > 50:
                df = df.head(49)
                rows = 50
            if cols > 10:
                df = df.iloc[:, :10]
                cols = 10
            
            col_widths = [min(max_col_width, max(len(str(col)) * 8 + cell_padding * 2, 60)) for col in df.columns]
            total_width = sum(col_widths) + cols + 1  # +1 for borders
            total_height = rows * row_height + rows + 1
            
            # Create image
            img = Image.new('RGB', (total_width, total_height), color=bg_color)
            draw = ImageDraw.Draw(img)
            
            # Draw table
            y = 0
            
            # Header row
            x = 0
            for i, col in enumerate(df.columns):
                # Draw header cell background
                draw.rectangle([x, y, x + col_widths[i], y + row_height], fill=header_bg, outline=border_color)
                # Draw header text
                text_x = x + cell_padding
                text_y = y + (row_height - 12) // 2
                draw.text((text_x, text_y), str(col), font=font, fill=text_color)
                x += col_widths[i] + 1
            y += row_height + 1
            
            # Data rows
            for _, row in df.iterrows():
                x = 0
                for i, (col, value) in enumerate(zip(df.columns, row)):
                    # Draw cell
                    draw.rectangle([x, y, x + col_widths[i], y + row_height], fill=bg_color, outline=border_color)
                    # Draw text
                    text = str(value) if pd.notna(value) else ''
                    if len(text) > 15:
                        text = text[:12] + '...'
                    text_x = x + cell_padding
                    text_y = y + (row_height - 12) // 2
                    draw.text((text_x, text_y), text, font=font, fill=text_color)
                    x += col_widths[i] + 1
                y += row_height + 1
            
            # Save image
            fmt = 'JPEG' if output_format.lower() in ['jpg', 'jpeg'] else 'PNG'
            img.save(output_path, fmt)
            
            logger.info(f"Successfully converted CSV to {fmt}: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"CSV to image conversion failed: {str(e)}")
            return False
    
    def convert_xlsx_to_image(self, input_path, output_path, output_format='png'):
        """Convert XLSX to image by rendering spreadsheet data as image."""
        try:
            import pandas as pd
            
            # Read first sheet of Excel file
            df = pd.read_excel(input_path, sheet_name=0)
            
            # Create temporary CSV and convert to image
            import tempfile
            temp_csv = os.path.join(tempfile.gettempdir(), f"xlsx_data_{uuid.uuid4()}.csv")
            try:
                df.to_csv(temp_csv, index=False)
                result = self.convert_csv_to_image(temp_csv, output_path, output_format)
                return result
            finally:
                try:
                    os.remove(temp_csv)
                except Exception:
                    pass
                    
        except Exception as e:
            logger.error(f"XLSX to image conversion failed: {str(e)}")
            return False
    
    def convert_csv_to_office(self, input_path, output_path, output_format):
        """Convert CSV to office formats using pandas and basic templates."""
        try:
            # Read CSV
            df = pd.read_csv(input_path)
            
            if output_format.lower() in ['xlsx', 'xls']:
                df.to_excel(output_path, index=False)
                return True
            elif output_format.lower() == 'html':
                html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>CSV Data</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 20px; }}
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; }}
    </style>
</head>
<body>
    <h1>CSV Data</h1>
    {df.to_html(index=False, table_id='data-table')}
</body>
</html>"""
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                return True
            elif output_format.lower() in ['pptx', 'ppt']:
                # Simple PowerPoint with table - requires python-pptx
                try:
                    from pptx import Presentation
                    from pptx.util import Inches
                    
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = "CSV Data"
                    
                    # Add table
                    rows, cols = len(df) + 1, len(df.columns)  # +1 for header
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(5)
                    
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Header row
                    for j, col_name in enumerate(df.columns):
                        table.cell(0, j).text = str(col_name)
                    
                    # Data rows
                    for i, (_, row) in enumerate(df.iterrows(), 1):
                        for j, value in enumerate(row):
                            table.cell(i, j).text = str(value)
                    
                    prs.save(output_path)
                    return True
                except ImportError:
                    logger.warning("python-pptx not available for CSV to PowerPoint conversion")
                    return False
            elif output_format.lower() in ['docx', 'doc']:
                # CSV to Word using python-docx
                try:
                    from docx import Document
                    
                    doc = Document()
                    doc.add_heading('CSV Data', 0)
                    
                    # Limit table size for Word compatibility
                    max_rows = 100
                    max_cols = 15
                    
                    display_df = df.copy()
                    if len(display_df) > max_rows:
                        doc.add_paragraph(f"Note: Showing first {max_rows} rows of {len(df)} total rows.")
                        display_df = display_df.head(max_rows)
                    
                    if len(display_df.columns) > max_cols:
                        doc.add_paragraph(f"Note: Showing first {max_cols} columns of {len(df.columns)} total columns.")
                        display_df = display_df.iloc[:, :max_cols]
                    
                    # Create table in Word document
                    if not display_df.empty:
                        rows = len(display_df) + 1  # +1 for header
                        cols = len(display_df.columns)
                        table = doc.add_table(rows=rows, cols=cols)
                        table.style = 'Table Grid'
                        
                        # Add headers
                        for j, column_name in enumerate(display_df.columns):
                            cell = table.cell(0, j)
                            cell.text = str(column_name)
                            # Make header bold
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                        
                        # Add data rows
                        for i, (_, row) in enumerate(display_df.iterrows()):
                            for j, value in enumerate(row):
                                table.cell(i + 1, j).text = str(value) if pd.notna(value) else ""
                    
                    doc.save(output_path)
                    return True
                except ImportError:
                    logger.warning("python-docx not available for CSV to Word conversion")
                    return False
            
            return False
        except Exception as e:
            logger.error(f"CSV conversion failed: {str(e)}")
            return False
    
    def convert_file(self, input_path, input_ext, output_format, original_filename):
        """Main conversion method that tries multiple engines with comprehensive coverage."""
        output_filename = os.path.splitext(os.path.basename(input_path))[0] + '.' + output_format
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        conversion_key = f"{input_ext}_to_{output_format}"
        logger.info(f"Converting {conversion_key}: {input_path} -> {output_path}")
        
        # === PRIORITY 1: Python libraries for specific optimized cases ===
        
        # Images to PDF
        if input_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff'] and output_format == 'pdf':
            if self.convert_image_to_pdf(input_path, output_path):
                return output_path
        
        # HTML to PDF
        if input_ext == 'html' and output_format == 'pdf' and self.has_wkhtmltopdf:
            if self.convert_html_to_pdf(input_path, output_path):
                return output_path
        
        # DOCX to PPTX (Python-based since LibreOffice doesn't support this well)
        if input_ext == 'docx' and output_format == 'pptx':
            try:
                if self.convert_docx_to_pptx(input_path, output_path):
                    return output_path
                else:
                    # If DOCX conversion fails, don't try other methods for DOCX files
                    raise Exception("Cannot convert DOCX file - it may be corrupted, empty, or not a valid Word document")
            except Exception as docx_error:
                logger.error(f"DOCX to PPTX conversion failed: {str(docx_error)}")
                raise docx_error
        
        # TXT conversions (Python-based)
        if input_ext == 'txt' and output_format == 'html':
            if self.convert_txt_to_html(input_path, output_path):
                return output_path
        
        if input_ext == 'txt' and output_format == 'pptx':
            if self.convert_txt_to_pptx(input_path, output_path):
                return output_path
        
        # HTML conversions (Python-based)
        if input_ext == 'html' and output_format == 'txt':
            if self.convert_html_to_txt(input_path, output_path):
                return output_path
        
        if input_ext == 'html' and output_format == 'pptx':
            if self.convert_html_to_pptx(input_path, output_path):
                return output_path
        
        # Image format conversions
        if input_ext == 'png' and output_format in ['jpg', 'jpeg']:
            if self.convert_image_formats(input_path, output_path, 'png', 'jpeg'):
                return output_path
        
        if input_ext in ['jpg', 'jpeg'] and output_format == 'png':
            if self.convert_image_formats(input_path, output_path, 'jpeg', 'png'):
                return output_path
        
        # Images to Office/HTML/TXT conversions
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'docx':
            if self.convert_image_to_docx(input_path, output_path):
                return output_path
        
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'pptx':
            if self.convert_image_to_pptx(input_path, output_path):
                return output_path
        
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'html':
            if self.convert_image_to_html(input_path, output_path):
                return output_path
        
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'txt':
            if self.convert_image_to_txt(input_path, output_path):
                return output_path
        
        # Images to CSV/XLSX conversions (NEW VICE-VERSA)
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'csv':
            if self.convert_image_to_csv(input_path, output_path):
                return output_path
        
        if input_ext in ['png', 'jpg', 'jpeg'] and output_format == 'xlsx':
            if self.convert_image_to_xlsx(input_path, output_path):
                return output_path
        
        # TXT/HTML to image conversions
        if input_ext == 'txt' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_txt_to_image(input_path, output_path, output_format):
                return output_path
        
        if input_ext == 'html' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_html_to_image(input_path, output_path, output_format):
                return output_path
        
        # DOCX/PPTX to image conversions (VICE-VERSA)
        if input_ext == 'docx' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_docx_to_image(input_path, output_path, output_format):
                return output_path
        
        if input_ext == 'pptx' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_pptx_to_image(input_path, output_path, output_format):
                return output_path
        
        # CSV/XLSX to image conversions (VICE-VERSA)
        if input_ext == 'csv' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_csv_to_image(input_path, output_path, output_format):
                return output_path
        
        if input_ext == 'xlsx' and output_format in ['png', 'jpg', 'jpeg']:
            if self.convert_xlsx_to_image(input_path, output_path, output_format):
                return output_path
        
        # PDF to text-based formats (extract text)
        if input_ext == 'pdf' and output_format in ['txt', 'html', 'csv', 'xlsx', 'xls']:
            if self.convert_pdf_extract_text(input_path, output_path, output_format):
                return output_path
        
        # PDF to images
        if input_ext == 'pdf' and output_format in ['jpg', 'jpeg', 'png']:
            if self.convert_pdf_to_images(input_path, output_path, output_format):
                return output_path
        
        # PDF to Word documents (using python-docx)
        if input_ext == 'pdf' and output_format in ['docx', 'doc']:
            if self.convert_pdf_to_word(input_path, output_path, output_format):
                return output_path
        
        # PDF to PowerPoint (using python-pptx)
        if input_ext == 'pdf' and output_format in ['pptx', 'ppt']:
            if self.convert_pdf_to_pptx(input_path, output_path):
                return output_path
        
        # CSV conversions (including Word documents)
        if input_ext == 'csv' and output_format in ['xlsx', 'xls', 'html', 'pptx', 'ppt', 'docx', 'doc']:
            if self.convert_csv_to_office(input_path, output_path, output_format):
                return output_path
        
        # PowerPoint to Word conversions (using python-pptx and python-docx)
        if input_ext in ['pptx', 'ppt'] and output_format in ['docx', 'doc']:
            if self.convert_pptx_to_word(input_path, output_path, output_format):
                return output_path
        
        # Excel to Word conversions (using pandas and python-docx)
        if input_ext in ['xlsx', 'xls'] and output_format in ['docx', 'doc']:
            if self.convert_excel_to_word(input_path, output_path, output_format):
                return output_path
        
        # PowerPoint to PDF conversions (using python-pptx)
        if input_ext in ['pptx', 'ppt'] and output_format == 'pdf':
            if self.convert_pptx_to_pdf(input_path, output_path):
                return output_path
        
        # === PRIORITY 2: LibreOffice for office document conversions ===
        
        if self.has_libreoffice:
            # Check if LibreOffice can handle this conversion (using our curated list)
            conversion_key = f"{input_ext}_to_{output_format}"
            capabilities = self.get_capabilities()
            
            if conversion_key in capabilities.get('libreoffice', []):
                
                try:
                    logger.info(f"Attempting LibreOffice conversion: {input_ext} -> {output_format}")
                    self.convert_with_libreoffice(input_path, output_format, input_ext)
                    
                    # LibreOffice creates files based on input filename
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    # Remove UUID prefix if present
                    if len(base_name) > 36 and base_name[8] == '-' and base_name[13] == '-':
                        actual_base_name = '_'.join(base_name.split('_')[1:])  # Remove UUID prefix
                    else:
                        actual_base_name = base_name
                    
                    # Try multiple possible output filenames
                    possible_outputs = [
                        os.path.join(OUTPUT_FOLDER, f"{base_name}.{output_format}"),
                        os.path.join(OUTPUT_FOLDER, f"{actual_base_name}.{output_format}"),
                        os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(original_filename)[0]}.{output_format}")
                    ]
                    
                    for possible_output in possible_outputs:
                        if os.path.exists(possible_output) and os.path.getsize(possible_output) > 0:
                            if possible_output != output_path:
                                try:
                                    os.rename(possible_output, output_path)
                                except OSError:
                                    # If rename fails, copy the file
                                    shutil.copy2(possible_output, output_path)
                                    os.remove(possible_output)
                            logger.info(f"LibreOffice conversion successful: {output_path}")
                            return output_path
                    
                    # If specific filename not found, look for any recent file with correct extension
                    for file in os.listdir(OUTPUT_FOLDER):
                        file_path = os.path.join(OUTPUT_FOLDER, file)
                        if (file.endswith(f'.{output_format}') and 
                            os.path.getmtime(file_path) > time.time() - 120 and  # Created in last 2 minutes
                            os.path.getsize(file_path) > 0):
                            
                            try:
                                os.rename(file_path, output_path)
                                logger.info(f"LibreOffice conversion successful (found recent file): {output_path}")
                                return output_path
                            except OSError:
                                shutil.copy2(file_path, output_path)
                                os.remove(file_path)
                                logger.info(f"LibreOffice conversion successful (copied recent file): {output_path}")
                                return output_path
                    
                    # If we get here, LibreOffice ran but didn't produce expected output
                    raise Exception(f"LibreOffice completed but output file not found")
                    
                except Exception as e:
                    logger.warning(f"LibreOffice conversion failed: {e}")
                    # Continue to next method instead of failing immediately
        
        # === PRIORITY 3: CloudConvert as fallback for unsupported conversions ===
        
        if CLOUDCONVERT_API_KEY:
            try:
                logger.info(f"Attempting CloudConvert conversion: {input_ext} -> {output_format}")
                converted_content = self.convert_with_cloudconvert(input_path, output_format)
                with open(output_path, 'wb') as f:
                    f.write(converted_content)
                logger.info(f"CloudConvert conversion successful: {output_path}")
                return output_path
            except Exception as e:
                logger.warning(f"CloudConvert conversion failed: {e}")
        
        # === If all methods fail ===
        
        available_methods = []
        if self.has_libreoffice:
            available_methods.append("LibreOffice")
        if self.has_wkhtmltopdf:
            available_methods.append("wkhtmltopdf")
        if CLOUDCONVERT_API_KEY:
            available_methods.append("CloudConvert API")
        
        error_msg = f"No conversion method available for {input_ext} to {output_format}."
        if available_methods:
            error_msg += f" Available conversion engines: {', '.join(available_methods)}."
        else:
            error_msg += " No conversion engines are properly configured."
        
        raise Exception(error_msg)

# Initialize conversion engine
conversion_engine = ConversionEngine()

@app.route('/')
def index():
    capabilities = conversion_engine.get_capabilities()
    return render_template('hybrid_index.html', 
                         capabilities=capabilities,
                         has_libreoffice=conversion_engine.has_libreoffice,
                         has_wkhtmltopdf=conversion_engine.has_wkhtmltopdf,
                         has_cloudconvert=bool(CLOUDCONVERT_API_KEY))

@app.route('/healthz')
def healthz():
    return 'ok', 200

@app.route('/capabilities')
def get_capabilities():
    """API endpoint to get current conversion capabilities"""
    return jsonify(conversion_engine.get_capabilities())

@app.route('/manifest.json')
def manifest():
    return send_file('templates/manifest.json', mimetype='application/json')

@app.route('/sw.js')
def service_worker():
    sw_content = '''
const CACHE_NAME = 'hybrid-converter-v1';
const urlsToCache = [
    '/',
    '/manifest.json',
    '/capabilities',
    'https://cdnjs.cloudflare.com/ajax/libs/pdf-lib/1.17.1/pdf-lib.min.js',
    'https://cdnjs.cloudflare.com/ajax/libs/jspdf/2.5.1/jspdf.umd.min.js',
    'https://unpkg.com/mammoth@1.4.21/mammoth.browser.min.js',
    'https://cdnjs.cloudflare.com/ajax/libs/xlsx/0.18.5/xlsx.full.min.js'
];

self.addEventListener('install', (event) => {
    event.waitUntil(
        caches.open(CACHE_NAME)
            .then((cache) => cache.addAll(urlsToCache))
    );
});

self.addEventListener('fetch', (event) => {
    // Only cache GET requests; pass through others (like POST /convert)
    if (event.request.method !== 'GET') {
        return;
    }
    event.respondWith(
        caches.match(event.request)
            .then((response) => response || fetch(event.request))
    );
});
    '''
    return app.response_class(sw_content, mimetype='application/javascript')

@app.route('/convert', methods=['POST'])
def convert():
    """Handle file conversion with comprehensive engine support."""
    try:
        # Validate request
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
            
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
            
        # Secure filename and validate
        original_filename = secure_filename(file.filename)
        if not original_filename or not allowed_file(original_filename):
            return jsonify({
                "error": f"Unsupported file type. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
            }), 400

        output_format = request.form.get('output_format')
        if not output_format:
            return jsonify({"error": "Output format not specified"}), 400

        # Create secure filenames
        input_filename = f"{uuid.uuid4()}_{original_filename}"
        input_path = os.path.join(UPLOAD_FOLDER, input_filename)
        
        logger.info(f"Processing conversion: {original_filename} -> {output_format}")
        file.save(input_path)
        
        input_ext = original_filename.rsplit('.', 1)[-1].lower()
        
        try:
            # Use the comprehensive conversion engine
            output_path = conversion_engine.convert_file(input_path, input_ext, output_format, original_filename)
            
            # Clean up input file
            try:
                os.remove(input_path)
            except OSError:
                logger.warning(f"Could not remove input file: {input_path}")
            
            # Validate output file
            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                raise Exception("Conversion produced invalid output file")
            
            logger.info(f"Conversion successful: {original_filename} -> {output_format}")
            return send_file(output_path, as_attachment=True, 
                           download_name=f"{os.path.splitext(original_filename)[0]}.{output_format}")
            
        except Exception as e:
            logger.error(f"Conversion failed: {str(e)}")
            
            # Determine what's missing and provide helpful error message
            error_msg = str(e)
            suggestions = []
            
            if not conversion_engine.has_libreoffice:
                suggestions.append("Install LibreOffice for office document conversions")
            if not conversion_engine.has_wkhtmltopdf:
                suggestions.append("Install wkhtmltopdf for HTML to PDF conversions")
            if not CLOUDCONVERT_API_KEY:
                suggestions.append("Configure CloudConvert API for additional format support")
            
            if suggestions:
                error_msg += f". Suggestions: {'; '.join(suggestions)}"
            
            return jsonify({"error": error_msg}), 400
        
    except Exception as e:
        logger.error(f"Unexpected error during conversion: {str(e)}")
        # Clean up input file if it exists
        try:
            if 'input_path' in locals() and os.path.exists(input_path):
                os.remove(input_path)
        except OSError:
            pass
            
        return jsonify({"error": "An unexpected error occurred during conversion."}), 500

# Add error handlers
@app.errorhandler(413)
def too_large(e):
    return jsonify({"error": "File too large. Maximum size is 50MB."}), 413

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {error}")
    return jsonify({"error": "Internal server error occurred."}), 500

if __name__ == '__main__':
    start_cleanup_thread()
    logger.info("Starting hybrid Flask application")
    logger.info(f"LibreOffice available: {conversion_engine.has_libreoffice}")
    logger.info(f"wkhtmltopdf available: {conversion_engine.has_wkhtmltopdf}")
    logger.info(f"CloudConvert API configured: {bool(CLOUDCONVERT_API_KEY)}")
    
    # Production vs Development configuration
    is_production = os.environ.get('FLASK_ENV') == 'production'
    port = int(os.environ.get('PORT', 5000))
    
    if is_production:
        # Production mode - let gunicorn handle it
        logger.info("Running in production mode")
    else:
        # Development mode
        logger.info("Running in development mode")
        app.run(debug=True, host='127.0.0.1', port=port)
